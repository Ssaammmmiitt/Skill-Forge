"""Extract plain text from PDF, Word, and PowerPoint uploads."""
import io
from typing import BinaryIO

from api.exceptions import ApiError

MAX_EXTRACT_CHARS = 50_000
ALLOWED_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt"}


def _truncate(text: str) -> str:
    text = text.strip()
    if len(text) > MAX_EXTRACT_CHARS:
        return text[:MAX_EXTRACT_CHARS] + "\n\n[... document truncated for processing ...]"
    return text


def extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    parts = []
    for page in reader.pages:
        parts.append(page.extract_text() or "")
    return _truncate("\n".join(parts))


def extract_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return _truncate("\n".join(parts))


def _pptx_text_frame_lines(text_frame) -> list[str]:
    lines = []
    for para in text_frame.paragraphs:
        runs = [run.text for run in para.runs if run.text]
        line = "".join(runs).strip() if runs else (para.text or "").strip()
        if line:
            lines.append(line)
    return lines


def _pptx_shape_text(shape) -> list[str]:
    """Collect text from a shape, including groups, tables, and placeholders."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    try:
        shape_type = shape.shape_type
    except Exception:
        return []

    parts: list[str] = []

    if shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            parts.extend(_pptx_shape_text(child))
        return parts

    if shape_type == MSO_SHAPE_TYPE.TABLE:
        try:
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        except Exception:
            pass
        return parts

    try:
        if shape.has_text_frame:
            parts.extend(_pptx_text_frame_lines(shape.text_frame))
            if parts:
                return parts
    except Exception:
        pass

    try:
        if shape.text and shape.text.strip():
            parts.append(shape.text.strip())
    except Exception:
        pass

    return parts


def extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_bits: list[str] = []
        for shape in slide.shapes:
            slide_bits.extend(_pptx_shape_text(shape))

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_bits.append(f"[Speaker notes] {notes}")

        # De-duplicate consecutive identical lines while preserving order
        deduped: list[str] = []
        for line in slide_bits:
            if not deduped or deduped[-1] != line:
                deduped.append(line)

        if deduped:
            parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(deduped))
    return _truncate("\n\n".join(parts))


def extract_text(filename: str, data: bytes) -> str:
    if not data:
        raise ApiError("Uploaded file is empty", 400)

    ext = ""
    if "." in filename:
        ext = filename.rsplit(".", 1)[-1].lower()
        ext = f".{ext}"

    if ext not in ALLOWED_EXTENSIONS:
        raise ApiError(
            "Unsupported file type. Upload PDF (.pdf), Word (.docx), or PowerPoint (.pptx).",
            400,
        )

    try:
        if ext == ".pdf":
            text = extract_pdf(data)
        elif ext in (".docx", ".doc"):
            if ext == ".doc":
                raise ApiError(
                    "Legacy .doc files are not supported. Save as .docx and upload again.",
                    400,
                )
            text = extract_docx(data)
        elif ext in (".pptx", ".ppt"):
            if ext == ".ppt":
                raise ApiError(
                    "Legacy .ppt files are not supported. Save as .pptx and upload again.",
                    400,
                )
            text = extract_pptx(data)
        else:
            raise ApiError("Unsupported file type", 400)
    except ApiError:
        raise
    except Exception as exc:
        raise ApiError(f"Could not read file: {exc}", 400) from exc

    stripped = (text or "").strip()
    if len(stripped) < 20:
        hint = ""
        if ext == ".pptx":
            hint = (
                " Slides may be image-only—add speaker notes in PowerPoint, "
                "or export to PDF with selectable text."
            )
        elif ext == ".pdf":
            hint = " The PDF may be scanned images—use a text-based PDF or add OCR."
        raise ApiError(
            f"Could not extract enough text from this file ({len(stripped)} characters).{hint}",
            400,
        )
    return text
