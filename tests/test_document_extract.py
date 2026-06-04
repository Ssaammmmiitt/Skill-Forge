"""Tests for document text extraction (especially PPTX)."""
import io

import pytest

pytest.importorskip("pptx")

from pptx import Presentation
from pptx.util import Inches

from engine.document_extract import extract_pptx, extract_text


def _minimal_pptx_bytes(*slide_texts: str) -> bytes:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for text in slide_texts:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        box.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_extract_pptx_plain_text():
    data = _minimal_pptx_bytes("Hello slide", "Second slide")
    text = extract_pptx(data)
    assert "Hello slide" in text
    assert "Second slide" in text
    assert "--- Slide 1 ---" in text
    assert "--- Slide 2 ---" in text


def test_extract_text_pptx_extension():
    data = _minimal_pptx_bytes("Enough characters here for the minimum length check.")
    text = extract_text("lecture.pptx", data)
    assert "Enough characters" in text
